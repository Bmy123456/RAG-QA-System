from pptx import Presentation


def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text.strip())
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        slide_texts.append(row_text)
        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)
