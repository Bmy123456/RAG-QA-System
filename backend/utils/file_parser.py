"""
文件解析层：将不同格式的原始文件转为统一的 ParseResult。
支持格式：PDF, Word, Excel, PPT, TXT, Markdown, HTML(网页), 图片(OCR), 邮件(EML)
"""

from __future__ import annotations

import email
import mimetypes
from dataclasses import dataclass, field
from pathlib import Path

# ---------------------------------------------------------------------------
# 数据类
# ---------------------------------------------------------------------------

@dataclass
class PageContent:
    """单页内容（PDF / PPT 等分页格式使用）"""
    page_number: int
    text: str


@dataclass
class ParseResult:
    """统一的解析输出"""
    text: str                                    # 全文纯文本
    pages: list[PageContent] | None = None       # 按页拆分（PDF/PPT 有）
    file_type: str = ""                          # pdf / docx / xlsx / …
    metadata: dict = field(default_factory=dict) # 额外元数据（标题、作者等）


# ---------------------------------------------------------------------------
# 格式 → 解析函数 的注册表
# ---------------------------------------------------------------------------

_PARSERS: dict[str, str] = {
    ".pdf":  "parse_pdf",
    ".docx": "parse_docx",
    ".doc":  "parse_docx",
    ".xlsx": "parse_xlsx",
    ".xls":  "parse_xlsx",
    ".pptx": "parse_pptx",
    ".ppt":  "parse_pptx",
    ".txt":  "parse_txt",
    ".md":   "parse_markdown",
    ".html": "parse_html",
    ".htm":  "parse_html",
    ".eml":  "parse_eml",
    ".png":  "parse_image",
    ".jpg":  "parse_image",
    ".jpeg": "parse_image",
    ".bmp":  "parse_image",
    ".tiff": "parse_image",
    ".tif":  "parse_image",
}


# ---------------------------------------------------------------------------
# 入口
# ---------------------------------------------------------------------------

def parse_file(file_path: str) -> ParseResult:
    """根据文件扩展名分发到对应解析器。

    参数:
        file_path: 文件路径（本地文件）

    返回:
        ParseResult

    异常:
        ValueError: 不支持的文件格式
        FileNotFoundError: 文件不存在
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    suffix = path.suffix.lower()
    if suffix not in _PARSERS:
        raise ValueError(f"不支持的文件格式: {suffix}（文件: {file_path}）")

    func_name = _PARSERS[suffix]
    func = globals()[func_name]
    return func(str(path))


# ---------------------------------------------------------------------------
# PDF — PyMuPDF (fitz)
# ---------------------------------------------------------------------------

def parse_pdf(path: str) -> ParseResult:
    import fitz  # PyMuPDF

    doc = fitz.open(path)
    pages: list[PageContent] = []
    full_text_parts: list[str] = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text("text")
        pages.append(PageContent(page_number=page_num + 1, text=text))
        full_text_parts.append(text)

    metadata = doc.metadata or {}
    doc.close()

    return ParseResult(
        text="\n\n".join(full_text_parts),
        pages=pages,
        file_type="pdf",
        metadata={"title": metadata.get("title"), "author": metadata.get("author")},
    )


# ---------------------------------------------------------------------------
# Word — python-docx
# ---------------------------------------------------------------------------

def parse_docx(path: str) -> ParseResult:
    from docx import Document

    doc = Document(path)
    paragraphs: list[str] = []
    for para in doc.paragraphs:
        paragraphs.append(para.text)

    # 同时提取表格内容
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))

    return ParseResult(
        text="\n\n".join(paragraphs),
        file_type="docx",
        metadata={},
    )


# ---------------------------------------------------------------------------
# Excel — openpyxl
# ---------------------------------------------------------------------------

def parse_xlsx(path: str) -> ParseResult:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"[Sheet: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))

    wb.close()
    return ParseResult(
        text="\n".join(parts),
        file_type="xlsx",
        metadata={"sheets": wb.sheetnames},
    )


# ---------------------------------------------------------------------------
# PPT — python-pptx
# ---------------------------------------------------------------------------

def parse_pptx(path: str) -> ParseResult:
    from pptx import Presentation

    prs = Presentation(path)
    pages: list[PageContent] = []
    full_text_parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        page_text = "\n".join(texts)
        pages.append(PageContent(page_number=slide_num, text=page_text))
        full_text_parts.append(page_text)

    return ParseResult(
        text="\n\n".join(full_text_parts),
        pages=pages,
        file_type="pptx",
        metadata={},
    )


# ---------------------------------------------------------------------------
# TXT — 直接读取
# ---------------------------------------------------------------------------

def parse_txt(path: str) -> ParseResult:
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    return ParseResult(text=text, file_type="txt", metadata={})


# ---------------------------------------------------------------------------
# Markdown — 直接读取（保留原始文本）
# ---------------------------------------------------------------------------

def parse_markdown(path: str) -> ParseResult:
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    return ParseResult(text=text, file_type="md", metadata={})


# ---------------------------------------------------------------------------
# HTML — trafilatura
# ---------------------------------------------------------------------------

def parse_html(path_or_url: str) -> ParseResult:
    import trafilatura

    # 判断是本地文件还是 URL
    if path_or_url.startswith(("http://", "https://")):
        downloaded = trafilatura.fetch_url(path_or_url)
        text = trafilatura.extract(downloaded, include_comments=False, include_tables=True) or ""
    else:
        raw = Path(path_or_url).read_text(encoding="utf-8", errors="replace")
        text = trafilatura.extract(raw, include_comments=False, include_tables=True) or ""

    return ParseResult(text=text, file_type="html", metadata={})


# ---------------------------------------------------------------------------
# 图片 OCR — PaddleOCR
# ---------------------------------------------------------------------------

def parse_image(path: str) -> ParseResult:
    from paddleocr import PaddleOCR

    ocr = PaddleOCR(use_angle_cls=True, lang="ch", show_log=False)
    result = ocr.ocr(path, cls=True)

    lines: list[str] = []
    if result and result[0]:
        for item in result[0]:
            # item: [box, (text, confidence)]
            text = item[1][0] if isinstance(item[1], (list, tuple)) else str(item[1])
            lines.append(text)

    return ParseResult(
        text="\n".join(lines),
        file_type="image",
        metadata={},
    )


# ---------------------------------------------------------------------------
# 邮件 (EML) — 标准库 email + trafilatura 提取 HTML 正文
# ---------------------------------------------------------------------------

def parse_eml(path: str) -> ParseResult:
    import trafilatura

    raw_bytes = Path(path).read_bytes()
    msg = email.message_from_bytes(raw_bytes)

    subject = msg.get("Subject", "")
    from_addr = msg.get("From", "")
    to_addr = msg.get("To", "")
    date = msg.get("Date", "")

    body_text = ""

    if msg.is_multipart():
        for part in msg.walk():
            content_type = part.get_content_type()
            if content_type == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    body_text = payload.decode(charset, errors="replace")
                    break
            elif content_type == "text/html" and not body_text:
                payload = part.get_payload(decode=True)
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    html = payload.decode(charset, errors="replace")
                    body_text = trafilatura.extract(html) or ""
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            charset = msg.get_content_charset() or "utf-8"
            raw_text = payload.decode(charset, errors="replace")
            if msg.get_content_type() == "text/html":
                body_text = trafilatura.extract(raw_text) or ""
            else:
                body_text = raw_text

    header = f"From: {from_addr}\nTo: {to_addr}\nDate: {date}\nSubject: {subject}\n"
    full_text = header + "\n" + body_text

    return ParseResult(
        text=full_text.strip(),
        file_type="eml",
        metadata={"subject": subject, "from": from_addr, "to": to_addr, "date": date},
    )
